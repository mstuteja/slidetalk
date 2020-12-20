#python pptx library to control and manipulate powerpoint using python
from pptx import Presentation

#Slide textbox identification tool
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Function to extract content data from user's powerpoint presentation

def extractText(path):
    prs = Presentation(path)
    allData = []
    slideCount = 0
    for slide in prs.slides:
        slideCount=slideCount+1
        slideData=[]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    slideData.append(run.text)
        allData.append(slideData)
    return allData

# Function to sort and categorize content data and return headings and subheadings for further computation

def getHeader(data):
    resData = []
    for slide in data:
        headers = []
        tempCount = 0
        if(len(slide)==0):
            headers.append("NO HEADER")
        else:
            for text in slide:
                if(tempCount<3):
                    headers.append(text)
                    tempCount = tempCount + 1
        resData.append(headers)
    return resData
